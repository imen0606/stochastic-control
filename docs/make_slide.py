#!/usr/bin/env python3
"""Generate the fundraising slide as .pptx (opens in Keynote)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Colors
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x4E, 0x9A, 0xF5)
LIGHT_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)
MID_GRAY = RGBColor(0x66, 0x66, 0x77)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFC)
STAT_BG = RGBColor(0x1A, 0x1A, 0x2E)
RED_ACCENT = RGBColor(0xE8, 0x4D, 0x4D)
AMBER = RGBColor(0xE8, 0x9D, 0x2D)
GREEN = RGBColor(0x34, 0xA8, 0x53)

bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE


def add_text(sl, left, top, width, height, text, font_size=14,
             bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    txBox = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rect(sl, left, top, width, height, fill_color, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = sl.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if rounded:
        shape.adjustments[0] = 0.04
    return shape


# ══════════════════════════════════════════════
# TITLE
# ══════════════════════════════════════════════

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(0.3), Inches(0.05), Inches(0.45))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text(slide, 0.85, 0.25, 11, 0.55,
         "Why Finance", font_size=32, bold=True, color=BLACK)

add_text(slide, 0.85, 0.72, 11, 0.3,
         "The highest-value industry for AI agents \u2014 and the one least equipped to train them",
         font_size=13, color=MID_GRAY)


# ══════════════════════════════════════════════
# ROW 1 — MARKET STATS (4 cards)
# ══════════════════════════════════════════════

stats = [
    ("$700\u2013800B", "projected banking cost\nsavings from agentic AI", "McKinsey, 2026"),
    ("$1.5\u20132B", "projected annual AI value\nat JPMorgan alone", "450+ use cases"),
    ("44%", "of finance teams deploying\nagentic AI in 2026", "Wolters Kluwer survey"),
    ("160+", "agentic AI use cases from\ntop 50 banks in 2025", "Evident AI Index"),
]

stat_w = 2.85
stat_gap = 0.2
sx = 0.6

for i, (big_num, desc, sub) in enumerate(stats):
    x = sx + i * (stat_w + stat_gap)
    y = 1.1

    add_rect(slide, x, y, stat_w, 1.05, STAT_BG)

    add_text(slide, x + 0.18, y + 0.05, stat_w - 0.36, 0.38,
             big_num, font_size=24, bold=True, color=ACCENT)

    add_text(slide, x + 0.18, y + 0.4, stat_w - 0.36, 0.35,
             desc, font_size=9.5, color=WHITE)

    add_text(slide, x + 0.18, y + 0.8, stat_w - 0.36, 0.2,
             sub, font_size=8, color=RGBColor(0x88, 0x88, 0x99))


# ══════════════════════════════════════════════
# ROW 2 — THE GAP (full width, compact)
# ══════════════════════════════════════════════

add_text(slide, 0.6, 2.3, 6, 0.3,
         "The Gap", font_size=17, bold=True, color=RED_ACCENT)

add_rect(slide, 0.6, 2.6, 7.3, 0.95, RGBColor(0xFD, 0xF0, 0xF0))

add_text(slide, 0.85, 2.65, 6.8, 0.35,
         "The industry is deploying AI agents fast \u2014 but has no way to train or evaluate them for sequential financial decisions.",
         font_size=11.5, bold=True, color=DARK_GRAY)

add_text(slide, 0.85, 2.98, 6.8, 0.5,
         "Current benchmarks test static reasoning. Financial decisions are sequential: "
         "act under uncertainty, adapt over time, manage risk across steps.",
         font_size=10, color=MID_GRAY)

# Callout stat
add_rect(slide, 8.15, 2.6, 4.55, 0.95, RGBColor(0xFB, 0xE4, 0xE4))

add_text(slide, 8.35, 2.68, 4.15, 0.4,
         "95%", font_size=32, bold=True, color=RED_ACCENT)

add_text(slide, 8.35, 3.05, 4.15, 0.45,
         "of enterprise AI pilots fail to deliver measurable "
         "financial impact.  \u2014 MIT / BCG, 2025",
         font_size=10.5, color=DARK_GRAY)


# ══════════════════════════════════════════════
# ROW 3 — WHAT WE'RE BUILDING (3 columns)
# ══════════════════════════════════════════════

add_text(slide, 0.6, 3.72, 12, 0.3,
         "What We\u2019re Building", font_size=17, bold=True, color=ACCENT)

add_text(slide, 0.6, 3.98, 12, 0.25,
         "RL environments for financial decision-making \u2014 with provably optimal benchmarks to train and evaluate against.",
         font_size=11.5, color=DARK_GRAY)

cols = [
    ("Agentic Environments",
     "AI agents interact with markets over multiple steps \u2014 "
     "observing, deciding, adapting. Not static Q&A."),
    ("Verifiable by Design",
     "Every environment has a known optimal strategy. "
     "We measure what the agent learned, not just one backtest."),
    ("Broad Financial Coverage",
     "Trading, portfolio construction, execution, risk management \u2014 "
     "anywhere decisions are sequential."),
]

col_w = 3.85
col_gap = 0.28
col_x = 0.6
col_y = 4.3

for i, (title, desc) in enumerate(cols):
    x = col_x + i * (col_w + col_gap)

    add_rect(slide, x, col_y, col_w, 1.05, LIGHT_BG)

    add_text(slide, x + 0.2, col_y + 0.08, col_w - 0.4, 0.25,
             title, font_size=12.5, bold=True, color=BLACK)

    add_text(slide, x + 0.2, col_y + 0.35, col_w - 0.4, 0.65,
             desc, font_size=10, color=MID_GRAY)


# ══════════════════════════════════════════════
# ROW 4 — TWO BOXES: AI ACT + BEYOND FINANCE
# ══════════════════════════════════════════════

# AI Act box (left)
add_rect(slide, 0.6, 5.55, 6.15, 1.3, RGBColor(0xFF, 0xF8, 0xEB))

add_text(slide, 0.85, 5.6, 5.7, 0.25,
         "EU AI Act \u2014 High-Risk Deadline: August 2026", font_size=11, bold=True, color=AMBER)

add_text(slide, 0.85, 5.85, 5.7, 0.9,
         "Financial AI classified high-risk. "
         "Auditability, accuracy, and risk management now mandatory. "
         "Fines up to 7% of global turnover.",
         font_size=10.5, color=DARK_GRAY)

# Beyond Finance box (right)
add_rect(slide, 6.98, 5.55, 5.72, 1.3, RGBColor(0xEE, 0xF8, 0xEE))

add_text(slide, 7.2, 5.6, 5.3, 0.25,
         "Beyond Finance", font_size=11, bold=True, color=GREEN)

add_text(slide, 7.2, 5.85, 5.3, 0.9,
         "The same sequential decision-making applies broadly: "
         "energy dispatch, supply chain, healthcare treatment protocols. "
         "Finance first \u2014 the capability transfers.",
         font_size=10.5, color=DARK_GRAY)


# ══════════════════════════════════════════════
# BOTTOM BAR
# ══════════════════════════════════════════════

bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = BLACK
bottom_bar.line.fill.background()

add_text(slide, 0.6, 7.08, 12.1, 0.35,
         "Finance is the starting point \u2014 it has exact optimal solutions.  The capability transfers.",
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# Save
output_path = "docs/finance_ai_training_ground.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
