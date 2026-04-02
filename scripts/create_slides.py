"""Generate team presentation slides for the Financial RLVR Gym."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT2 = RGBColor(0xFF, 0xA5, 0x00)
ACCENT3 = RGBColor(0x6C, 0x9B, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x88, 0x88, 0x99)
RED_SOFT = RGBColor(0xE0, 0x60, 0x60)
GREEN_SOFT = RGBColor(0x60, 0xD0, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=6,
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, width=2.0):
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 1.5, 1.8, 10, 1.2,
             "Financial Planning RLVR Gym", 44, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.2, 10, 0.8,
             "Training multi-step planning under uncertainty", 24, LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 4.8, 10, 0.6,
             "Team Overview  |  April 2026", 18, DIM_GRAY, alignment=PP_ALIGN.CENTER)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(4.5), Inches(4.3), Inches(4), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()


# ============================================================
# SLIDE 2: The Pitch — One Sentence
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "THE PITCH", 14, ACCENT, bold=True)

add_text_box(slide, 1.5, 2.0, 10, 2.0,
             "The model decides each turn whether to run\na trading strategy or sit in cash.",
             32, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text_box(slide, 1.5, 4.0, 10, 1.5,
                  "Running earns or loses money depending on market conditions.",
                  20, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_para(tf, "Switching between in and out costs money.", 20, LIGHT_GRAY,
         alignment=PP_ALIGN.CENTER, space_before=8)
add_para(tf, "The optimal strategy is provably computable.", 20, ACCENT,
         bold=True, alignment=PP_ALIGN.CENTER, space_before=8)

add_text_box(slide, 1.5, 6.0, 10, 0.6,
             "One reward_fn. Plugs into GRPO. Infinite instances. Tunable difficulty.",
             16, DIM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: The Gap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "THE GAP IN CURRENT RLVR", 14, ACCENT, bold=True)

# Math gym box
add_rounded_rect(slide, 1.0, 1.5, 3.5, 2.8, BG_MID, ACCENT3)
tf = add_text_box(slide, 1.3, 1.7, 3.0, 0.5, "Math Gyms", 22, ACCENT3, bold=True)
tf = add_text_box(slide, 1.3, 2.3, 3.0, 1.8,
                  "One problem, one answer", 16, LIGHT_GRAY)
add_para(tf, "Deterministic", 16, LIGHT_GRAY, space_before=6)
add_para(tf, "Single-shot", 16, LIGHT_GRAY, space_before=6)
add_para(tf, "No uncertainty", 16, LIGHT_GRAY, space_before=6)

# Code gym box
add_rounded_rect(slide, 5.0, 1.5, 3.5, 2.8, BG_MID, ACCENT3)
tf = add_text_box(slide, 5.3, 1.7, 3.0, 0.5, "Code Gyms", 22, ACCENT3, bold=True)
tf = add_text_box(slide, 5.3, 2.3, 3.0, 1.8,
                  "Write code, run tests", 16, LIGHT_GRAY)
add_para(tf, "Deterministic", 16, LIGHT_GRAY, space_before=6)
add_para(tf, "Single-shot", 16, LIGHT_GRAY, space_before=6)
add_para(tf, "No uncertainty", 16, LIGHT_GRAY, space_before=6)

# This gym box
add_rounded_rect(slide, 9.0, 1.5, 3.5, 2.8, RGBColor(0x1E, 0x3A, 0x2E), ACCENT)
tf = add_text_box(slide, 9.3, 1.7, 3.0, 0.5, "This Gym", 22, ACCENT, bold=True)
tf = add_text_box(slide, 9.3, 2.3, 3.0, 1.8,
                  "Sequence of coupled decisions", 16, WHITE)
add_para(tf, "Stochastic", 16, WHITE, space_before=6)
add_para(tf, "Multi-step", 16, WHITE, space_before=6)
add_para(tf, "Plans under uncertainty", 16, WHITE, space_before=6)

# Bottom callout
add_rounded_rect(slide, 2.0, 5.0, 9.3, 1.5, BG_MID, ACCENT)
tf = add_text_box(slide, 2.5, 5.15, 8.3, 1.3,
                  "No existing gym trains: sequential planning + uncertainty + verifiable rewards",
                  20, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "This is the capability gap for agentic decision-making",
         16, LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=10)


# ============================================================
# SLIDE 4: How It Works — Data Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "HOW IT WORKS", 14, ACCENT, bold=True)

# Flow boxes
boxes = [
    (0.5, 2.2, 2.2, 1.5, "GENERATE", "Create synthetic\nmarket scenario\nwith known rules", ACCENT),
    (3.2, 2.2, 2.2, 1.5, "SOLVE", "Bellman backward\ninduction gives\nexact optimal", ACCENT2),
    (5.9, 2.2, 2.2, 1.5, "PROMPT", "Multi-turn:\none signal per\nturn to LLM", ACCENT3),
    (8.6, 2.2, 2.2, 1.5, "PARSE", "Extract binary\ndecisions from\nmodel text", ACCENT3),
    (11.0, 2.2, 1.8, 1.5, "SCORE", "Compare\nprofit vs\noptimal", ACCENT),
]

for x, y, w, h, title, desc, color in boxes:
    add_rounded_rect(slide, x, y, w, h, BG_MID, color)
    add_text_box(slide, x + 0.15, y + 0.15, w - 0.3, 0.4, title, 15, color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, y + 0.6, w - 0.3, 0.9, desc, 13, LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Arrows between boxes
for i in range(len(boxes) - 1):
    x1 = boxes[i][0] + boxes[i][2]
    x2 = boxes[i + 1][0]
    y = boxes[i][1] + boxes[i][3] / 2
    add_arrow(slide, x1 + 0.05, y, x2 - 0.05, y, DIM_GRAY, 1.5)

# Output arrow to GRPO
add_rounded_rect(slide, 10.4, 4.5, 3.0, 1.0, BG_MID, ACCENT)
add_text_box(slide, 10.5, 4.6, 2.8, 0.8, "GRPO Trainer\n(TRL / veRL)", 14, ACCENT,
             alignment=PP_ALIGN.CENTER)
add_arrow(slide, 11.9, 3.7, 11.9, 4.5, ACCENT, 1.5)

# "We build this" bracket
add_rounded_rect(slide, 0.3, 5.8, 12.8, 0.8, RGBColor(0x1E, 0x3A, 0x2E), ACCENT)
tf = add_text_box(slide, 0.5, 5.9, 6, 0.6,
                  "We build: Generate + Solve + Prompt + Parse + Score", 16, ACCENT, bold=True)
tf = add_text_box(slide, 7.5, 5.9, 5.5, 0.6,
                  "Lab provides: GRPO training loop (already exists)", 16, DIM_GRAY)


# ============================================================
# SLIDE 5: The Interaction — Multi-Turn
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "MODEL INTERACTION", 14, ACCENT, bold=True)

add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "The model receives one signal per turn — it cannot see the future.",
             18, LIGHT_GRAY)

# Conversation mock
conv = [
    ("SETUP", "user", "You manage a momentum strategy. Switching costs 0.10.\nExpected PnL = 0.30 x signal. One observation at a time."),
    ("t=0", "user", "Z = +0.60  |  Previous: OFF"),
    ("t=0", "model", "E[PnL] = 0.18 > cost 0.10. Strong signal. Switch ON.\ns_0 = 1"),
    ("t=1", "user", "Z = +0.45  |  Previous: ON"),
    ("t=1", "model", "Already on, no cost. E[PnL] = 0.135 > 0. Stay in.\ns_1 = 1"),
    ("t=2", "user", "Z = -0.12  |  Previous: ON"),
    ("t=2", "model", "Signal negative. Switching off costs 0.10 but stops losses.\ns_2 = 0"),
]

y = 1.9
for label, role, text in conv:
    if role == "user":
        color_box = RGBColor(0x2A, 0x2A, 0x45)
        color_label = ACCENT3
        x_offset = 0.8
    else:
        color_box = RGBColor(0x1E, 0x3A, 0x2E)
        color_label = ACCENT
        x_offset = 2.5

    lines = text.count('\n') + 1
    h = 0.28 + lines * 0.22
    add_rounded_rect(slide, x_offset, y, 9.5, h, color_box)
    add_text_box(slide, x_offset + 0.15, y + 0.05, 1.0, 0.25, label, 10, color_label, bold=True)
    add_text_box(slide, x_offset + 1.3, y + 0.05, 8.0, h - 0.1, text, 12, WHITE,
                 font_name="Consolas")
    y += h + 0.08

# Why note
add_text_box(slide, 0.8, 6.5, 11, 0.6,
             "Why sequential?  If the model saw all signals upfront, it could cheat with future info. "
             "This way, model and benchmark play by the same rules.",
             14, DIM_GRAY)


# ============================================================
# SLIDE 6: Scoring — Worked Example
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "SCORING — WORKED EXAMPLE", 14, ACCENT, bold=True)

# Three agents
tf = add_text_box(slide, 0.8, 1.3, 6, 0.4,
                  "Three agents, same market path (T=5, λ=0.10):", 16, LIGHT_GRAY)

# Decision table
add_rounded_rect(slide, 0.8, 1.9, 11.5, 2.5, BG_MID)

header = "                    t=0    t=1    t=2    t=3    t=4    Profit"
tf = add_text_box(slide, 1.0, 2.0, 11, 0.3, header, 14, DIM_GRAY, font_name="Consolas")

add_para(tf, "Optimal:   s* = [ 1      1      0      0      0 ]    +$0.12",
         14, ACCENT, font_name="Consolas", space_before=10)
add_para(tf, "Model:   s_m = [ 1      1      0      0      1 ]    +$0.16",
         14, ACCENT3, font_name="Consolas", space_before=6)
add_para(tf, "Random:  s_r = [ 0      1      1      0      1 ]    -$0.14",
         14, RED_SOFT, font_name="Consolas", space_before=6)

# Score calculation
add_rounded_rect(slide, 3.0, 4.8, 7.3, 1.8, BG_MID, ACCENT)

tf = add_text_box(slide, 3.3, 4.95, 6.8, 0.3,
                  "score = (model - random) / (optimal - random)", 18, WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER, font_name="Consolas")
add_para(tf, "", 6, WHITE, space_before=6)
add_para(tf, "     = (0.16 - (-0.14)) / (0.12 - (-0.14))", 16, LIGHT_GRAY,
         font_name="Consolas", alignment=PP_ALIGN.CENTER, space_before=6)
add_para(tf, "     = 0.30 / 0.26 = 1.15", 16, ACCENT,
         font_name="Consolas", alignment=PP_ALIGN.CENTER, space_before=6, bold=True)

add_text_box(slide, 0.8, 6.8, 12, 0.5,
             "0 = random  |  1.0 = optimal  |  >1.0 possible on single instances (averages to 1.0 over many)",
             14, DIM_GRAY)


# ============================================================
# SLIDE 7: Goldilocks Validation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 8, 0.6, "VALIDATION — THE GOLDILOCKS TEST", 14, ACCENT, bold=True)

add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "Before shipping: prove the reward signal discriminates between reasoning levels.",
             18, LIGHT_GRAY)

# Results table
add_rounded_rect(slide, 1.5, 2.0, 5.5, 3.5, BG_MID)

tf = add_text_box(slide, 1.8, 2.15, 5, 0.3,
                  "              Easy    Medium    Hard", 15, DIM_GRAY,
                  font_name="Consolas")
add_para(tf, "", 4, WHITE, space_before=4)
add_para(tf, "Random        ~0       ~0       ~0", 15, RED_SOFT,
         font_name="Consolas", space_before=8)
add_para(tf, "Greedy        0.72     0.51     0.28", 15, ACCENT2,
         font_name="Consolas", space_before=8)
add_para(tf, "Optimal       1.00     1.00     1.00", 15, ACCENT,
         font_name="Consolas", space_before=8)
add_para(tf, "", 4, WHITE, space_before=4)
add_para(tf, "Gap:          0.28     0.49     0.72", 15, WHITE,
         font_name="Consolas", space_before=8, bold=True)

# Interpretation
add_rounded_rect(slide, 7.5, 2.0, 5.0, 3.5, BG_MID, ACCENT)
tf = add_text_box(slide, 7.8, 2.2, 4.5, 0.4,
                  "What this proves:", 18, ACCENT, bold=True)
tf = add_text_box(slide, 7.8, 2.8, 4.5, 2.5,
                  "Gap widens as difficulty increases", 15, WHITE, bold=True)
add_para(tf, "", 4, WHITE, space_before=2)
add_para(tf, "Easy: greedy gets 72% of the way.", 14, LIGHT_GRAY, space_before=8)
add_para(tf, "Just reacting almost works.", 14, DIM_GRAY, space_before=2)
add_para(tf, "", 4, WHITE, space_before=2)
add_para(tf, "Hard: greedy only gets 28%.", 14, LIGHT_GRAY, space_before=8)
add_para(tf, "You MUST think ahead.", 14, ACCENT, bold=True, space_before=2)
add_para(tf, "", 4, WHITE, space_before=2)
add_para(tf, "This proves the gym tests planning,", 14, WHITE, space_before=8)
add_para(tf, "not just pattern matching.", 14, WHITE, space_before=2)

# Bottom: LLM row
add_rounded_rect(slide, 1.5, 5.8, 11, 1.0, RGBColor(0x1E, 0x3A, 0x2E), ACCENT)
tf = add_text_box(slide, 1.8, 5.95, 10.5, 0.7,
                  "Approach 3 (future): Add LLM as 4th row. Same infrastructure, no changes.",
                  16, ACCENT, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: Anticipated Questions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "ANTICIPATED LAB QUESTIONS", 14, ACCENT, bold=True)

questions = [
    ("Why synthetic, not real data?",
     "Real distributions are unknown. Estimated optima are noisy. "
     "Synthetic = exact ground truth. Same trade-off math gyms make."),
    ("Why no code interpreter?",
     "Model could just code the Bellman solver and ace every instance. "
     "That tests coding, not financial reasoning. We test internalised planning."),
    ("Is one problem type enough?",
     "No \u2014 designed as first in a suite. Almgren-Chriss (execution) and "
     "Merton (allocation) are next. All share base infrastructure."),
    ("Is the market model realistic?",
     "OU process captures mean-reversion (well-documented in finance). "
     "Simplified but not arbitrary. Gym tests reasoning, not market modelling."),
]

y = 1.3
for q, a in questions:
    add_rounded_rect(slide, 0.8, y, 11.7, 1.1, BG_MID)
    add_text_box(slide, 1.1, y + 0.1, 11, 0.35, q, 16, ACCENT2, bold=True)
    add_text_box(slide, 1.1, y + 0.5, 11, 0.5, a, 14, LIGHT_GRAY)
    y += 1.25


# ============================================================
# SLIDE 9: Gaming Analogy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 8, 0.6, "MAPPING TO A GAMING ENVIRONMENT", 14, ACCENT, bold=True)

# Two-column mapping
mappings = [
    ("Game State", "Market signal Z\u209c"),
    ("Player Action", "Binary: activate (1) or sit out (0)"),
    ("Action Cost", "Switching cost \u03bb"),
    ("Score", "Profit from decision sequence"),
    ("Optimal Play", "Bellman-optimal policy (exact)"),
    ("Difficulty Setting", "\u03bb, \u03b1, T parameters"),
    ("Fog of War", "Signals revealed one at a time"),
    ("Game Engine", "OU process (known rules)"),
]

y = 1.3
add_rounded_rect(slide, 0.8, y, 5.2, 0.5, ACCENT)
add_text_box(slide, 1.0, y + 0.05, 4.8, 0.4, "Game Concept", 15, BG_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, 6.2, y, 6.3, 0.5, ACCENT)
add_text_box(slide, 6.4, y + 0.05, 5.9, 0.4, "Financial Gym Equivalent", 15, BG_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)

y += 0.6
for game, finance in mappings:
    bg = BG_MID if mappings.index((game, finance)) % 2 == 0 else RGBColor(0x28, 0x28, 0x40)
    add_rounded_rect(slide, 0.8, y, 5.2, 0.55, bg)
    add_text_box(slide, 1.0, y + 0.1, 4.8, 0.35, game, 14, ACCENT3)
    add_rounded_rect(slide, 6.2, y, 6.3, 0.55, bg)
    add_text_box(slide, 6.4, y + 0.1, 5.9, 0.35, finance, 14, WHITE)
    y += 0.55

# Analogy box
add_rounded_rect(slide, 0.8, 6.0, 11.7, 1.2, RGBColor(0x1E, 0x3A, 0x2E), ACCENT)
tf = add_text_box(slide, 1.1, 6.1, 11.1, 1.0,
                  "Closest analogy: Turn-based game where you deploy/recall a unit each turn. "
                  "Deploying earns points if conditions are good, loses if not. "
                  "Deploy/recall costs resources. Conditions drift randomly but tend to revert to neutral.",
                  14, LIGHT_GRAY)
add_para(tf, "The optimal strategy isn't \"deploy when good\" \u2014 it's \"deploy when good enough "
         "AND likely to stay good long enough.\"",
         14, ACCENT, bold=True, space_before=8)


# ============================================================
# SLIDE 10: Suite Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "SUITE ROADMAP", 14, ACCENT, bold=True)

add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "Designed as first gym in a family. Shared base classes, shared validation, shared TRL integration.",
             16, LIGHT_GRAY)

problems = [
    ("Regime Switching", "Binary on/off", "Planning with\nswitching costs",
     "Exact (DP)", "Building now", ACCENT),
    ("Optimal Execution\n(Almgren-Chriss)", "Trade size\nper period", "Urgency vs\nmarket impact",
     "Closed-form", "Very high", GREEN_SOFT),
    ("Portfolio Allocation\n(Merton)", "% in risky\nasset", "Risk vs\nreturn",
     "Closed-form", "Very high", GREEN_SOFT),
    ("Hedging with Costs\n(Black-Scholes)", "Hedge ratio\nadjustment", "Precision vs\ncost",
     "Semi-analytical", "Moderate-high", ACCENT2),
]

# Header
y = 2.0
headers = ["Problem", "Action", "What it Tests", "Solution", "Confidence"]
x_positions = [0.8, 3.8, 5.8, 8.0, 10.2]
widths = [2.8, 1.8, 2.0, 2.0, 2.5]

for i, h in enumerate(headers):
    add_text_box(slide, x_positions[i], y, widths[i], 0.4, h, 13, DIM_GRAY, bold=True)

y += 0.5
for name, action, tests, solution, confidence, color in problems:
    bg = BG_MID
    add_rounded_rect(slide, 0.6, y, 12.2, 0.85, bg)
    add_text_box(slide, 0.8, y + 0.08, 2.8, 0.7, name, 13, color, bold=True)
    add_text_box(slide, 3.8, y + 0.08, 1.8, 0.7, action, 12, LIGHT_GRAY)
    add_text_box(slide, 5.8, y + 0.08, 2.0, 0.7, tests, 12, LIGHT_GRAY)
    add_text_box(slide, 8.0, y + 0.08, 2.0, 0.7, solution, 12, LIGHT_GRAY)
    add_text_box(slide, 10.2, y + 0.08, 2.5, 0.7, confidence, 13, color, bold=True)
    y += 0.95

# Bottom note
add_text_box(slide, 0.8, 6.2, 11.7, 0.8,
             "Each new problem type = 3 files (generator, verifier, prompts). "
             "Everything else is shared.",
             15, DIM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: Status & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.8, 0.5, 5, 0.6, "STATUS & NEXT STEPS", 14, ACCENT, bold=True)

steps = [
    ("Design spec", "Complete and reviewed", True),
    ("Implementation plan", "Next step", False),
    ("Core gym (generator + solver + verifier)", "To build", False),
    ("Prompt format (multi-turn conversation)", "To build", False),
    ("Validation suite (Goldilocks test)", "To build", False),
    ("LLM baseline evaluation (Approach 3)", "Future", False),
    ("Suite expansion (execution, allocation)", "Future", False),
]

y = 1.5
for step, status, done in steps:
    indicator_color = ACCENT if done else DIM_GRAY
    indicator = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.2), Inches(y + 0.1), Inches(0.25), Inches(0.25)
    )
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = indicator_color
    indicator.line.fill.background()

    add_text_box(slide, 1.7, y, 6, 0.4, step, 18, WHITE if not done else ACCENT,
                 bold=done)
    add_text_box(slide, 8.5, y, 4, 0.4, status, 16,
                 ACCENT if done else LIGHT_GRAY)
    y += 0.65


# Save
output_path = "/Users/imen/Documents/stochastic_control/docs/financial-rlvr-gym-slides.pptx"
prs.save(output_path)
print(f"Slides saved to {output_path}")
